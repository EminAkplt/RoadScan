# -*- coding: utf-8 -*-
"""RoadScan proje sunumunu (PPTX) üretir."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DOCS = os.path.join(ROOT, "docs")
FIGS = os.path.join(DOCS, "figs")
OUT = os.path.join(DOCS, "RoadScan_Sunum.pptx")

# Tema renkleri
BG     = RGBColor(0x0B, 0x0F, 0x17)   # koyu lacivert
PANEL  = RGBColor(0x14, 0x1D, 0x2E)
ACCENT = RGBColor(0x2D, 0xD4, 0xBF)   # teal
TXT    = RGBColor(0xE8, 0xEE, 0xF6)
MUTE   = RGBColor(0x9A, 0xA7, 0xBC)
KRIT   = RGBColor(0xEF, 0x44, 0x44)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def rect(slide, l, t, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def set_run(r, text, size, color, bold=False, italic=False, font="Calibri"):
    r.text = text; r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic; r.font.name = font


def content_slide(title, bullets, accent_tag="RoadScan"):
    s = prs.slides.add_slide(BLANK); bg(s)
    # üst şerit
    rect(s, 0, 0, 13.333, 0.18, ACCENT)
    # başlık
    tf = box(s, 0.6, 0.45, 12.1, 1.0)
    p = tf.paragraphs[0]; set_run(p.add_run(), title, 30, TXT, bold=True)
    # küçük etiket
    tf2 = box(s, 0.62, 1.45, 12, 0.4)
    set_run(tf2.paragraphs[0].add_run(), accent_tag, 13, ACCENT, bold=True)
    # gövde
    tf3 = box(s, 0.7, 2.0, 12.0, 5.0)
    for i, (lvl, text, *opt) in enumerate(bullets):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        p.level = lvl
        col = opt[0] if opt else TXT
        bold = len(opt) > 1 and opt[1]
        bullet = "•  " if lvl == 0 else "–  "
        set_run(p.add_run(), bullet + text, 19 if lvl == 0 else 16, col, bold=bold)
        p.space_after = Pt(7)
    return s


# ---------- 1. KAPAK ----------
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, 0, 0, 13.333, 0.25, ACCENT)
rect(s, 0, 7.25, 13.333, 0.25, ACCENT)
tf = box(s, 1.0, 2.0, 11.3, 2.2)
set_run(tf.paragraphs[0].add_run(), "🛣️ RoadScan", 54, ACCENT, bold=True)
p = tf.add_paragraph(); set_run(p.add_run(), "Akıllı Yol Bozukluğu Tespit Sistemi", 30, TXT, bold=True)
p = tf.add_paragraph(); set_run(p.add_run(), "Dashcam + Derin Öğrenme + Görüntü İşleme", 18, MUTE)
tf = box(s, 1.0, 5.3, 11.3, 1.6)
p = tf.paragraphs[0]; set_run(p.add_run(), "Mehmet Emin AKPOLAT", 20, TXT, bold=True)
p = tf.add_paragraph(); set_run(p.add_run(), "Fırat Üniversitesi · Teknoloji Fakültesi · Yazılım Mühendisliği", 15, MUTE)
p = tf.add_paragraph(); set_run(p.add_run(), "215542003@firat.edu.tr", 15, ACCENT)

# ---------- 2. PROBLEM & ÇÖZÜM ----------
content_slide("Problem ve Çözüm", [
    (0, "Problem: yoldaki çukur ve çatlaklar trafik güvenliğini tehdit eder; tespit manuel, yavaş, pahalı, kayıt dağınık.", MUTE),
    (0, "Çözüm: araç kamerası (dashcam) yolu gerçek zamanlı tarar.", TXT, True),
    (1, "Bozuklukları yapay zekâ ile tespit eder"),
    (1, "Konum + zaman + fotoğrafla eşleştirir"),
    (1, "Merkezi veritabanına raporlar, haritada gösterir"),
    (0, "Kurulum gerektirmez — tarayıcıda, internetsiz, cihaz üstünde çalışır.", ACCENT, True),
], "Giriş")

# ---------- 3. SİSTEM NASIL ÇALIŞIR ----------
content_slide("Sistem Nasıl Çalışır?", [
    (0, "Kamera / Video  →  Tespit Motoru (tarayıcı)", TXT, True),
    (1, "YOLO modeli + görüntü işleme ile bozukluk tespiti"),
    (0, "→  Backend API (Node/Express)  →  PostgreSQL", TXT, True),
    (1, "Tespit + konum + zaman + kırpılmış foto kaydı"),
    (0, "→  Yönetim Paneli (Leaflet harita)", TXT, True),
    (1, "Renk kodlu pinler, liste, filtre, istatistik"),
    (0, "3 ana modül: Tespit Motoru · Backend+Veritabanı · Yönetim Paneli", ACCENT, True),
], "Mimari")

# ---------- 4. GELİŞİM ----------
content_slide("Projenin Gelişimi: Klasik → YZ → Hibrit", [
    (0, "v1 — Klasik görüntü işleme (Canny kenar + kontur)", MUTE, True),
    (1, "'Çukur' kavramı yok; her kenarı bozukluk sandı → çok yanlış pozitif"),
    (0, "v2 — Yapay zekâ (YOLO)", TXT, True),
    (1, "Bozukluğun ne olduğunu öğrenen model → yanlış pozitif azaldı"),
    (0, "v3 — Hibrit (ikisi birlikte)", ACCENT, True),
    (1, "YOLO tespit eder, klasik görüntü işleme doğrular → en sağlam yapı"),
], "Yaklaşım")

# ---------- 5. MODEL ----------
content_slide("Yapay Zekâ / Model Kısmı", [
    (0, "Model: YOLOv8s — gerçek zamanlı, hafif, edge-dostu nesne tespiti", TXT, True),
    (0, "Transfer öğrenme: COCO ön-eğitimli ağırlıklardan ince ayar"),
    (0, "Veri: 3 açık set birleştirildi → ~46.000 görüntü, 4 sınıf", ACCENT, True),
    (1, "RDD2022 + BharatPotHole + IVCNZ (çukur 3 kaynaktan)"),
    (1, "Sınıflar: Boyuna / Enine / Timsah Çatlak + Çukur"),
    (0, "Eğitim: RTX 5060, 80 epoch, ~13 saat"),
    (0, "Dağıtım: ONNX → tarayıcıda, internetsiz (WebGPU/WASM)", ACCENT, True),
], "Model")

# ---------- 6. GÖRÜNTÜ İŞLEME ----------
content_slide("Görüntü İşleme Kısmı", [
    (0, "Ön işleme: letterbox 640×640 + normalize → tensör", TXT, True),
    (0, "Son işleme: çıktı ayrıştırma + NMS (çakışan kutuları eleme)", TXT, True),
    (0, "Hibrit doğrulama (klasik CV) — çöp Kritik'leri eler:", ACCENT, True),
    (1, "ROI: kutu üstteyse (gökyüzü/bina/ağaç) → reddet"),
    (1, "Doku/Gradyan: düz/tek renk yüzey (araç/duvar) → reddet"),
    (0, "Ayrıca: severity (boyut), kırpma, kare-arası IoU takibi", MUTE),
], "Görüntü İşleme")

# ---------- 7. AKILLI RAPORLAMA ----------
content_slide("Akıllı Raporlama", [
    (0, "Ekranda 3 kademe de kutulanır: Küçük / Orta / Kritik", TXT, True),
    (0, "Panele YALNIZCA Kritik gönderilir → veritabanı şişmez", KRIT, True),
    (0, "Kare-arası takip: aynı çukur videoda tek kayıt"),
    (0, "Tespit bölgesi kırpılıp panele yakın-çekim gider"),
    (0, "Tüm eşikler sabit — son kullanıcı ayar yapmaz", MUTE),
], "Raporlama")

# ---------- 8. SONUÇLAR (figürlü) ----------
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, 0, 0, 13.333, 0.18, ACCENT)
tf = box(s, 0.6, 0.4, 12, 1.0); set_run(tf.paragraphs[0].add_run(), "Deneysel Sonuçlar", 30, TXT, bold=True)
tf2 = box(s, 0.62, 1.35, 12, 0.4); set_run(tf2.paragraphs[0].add_run(), "Bulgular", 13, ACCENT, bold=True)
# metin
tf3 = box(s, 0.7, 1.95, 5.4, 5.0)
res = [("Genel mAP@50", "0.611", ACCENT), ("Timsah Çatlak", "0.681", TXT), ("Enine Çatlak", "0.603", TXT),
       ("Boyuna Çatlak", "0.602", TXT), ("Çukur", "0.559", TXT)]
for i, (k, v, c) in enumerate(res):
    p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
    set_run(p.add_run(), f"• {k}:  ", 18, TXT, bold=(i == 0))
    set_run(p.add_run(), v, 18, c, bold=True)
    p.space_after = Pt(10)
p = tf3.add_paragraph(); p.space_before = Pt(8)
set_run(p.add_run(), "Doğrulama: 7.226 görüntü · 80 epoch", 13, MUTE)
# figür
if os.path.exists(os.path.join(FIGS, "fig2.png")):
    s.shapes.add_picture(os.path.join(FIGS, "fig2.png"), Inches(6.4), Inches(2.1), width=Inches(6.4))
tf4 = box(s, 6.4, 6.5, 6.4, 0.5)
set_run(tf4.paragraphs[0].add_run(), "Sınıf bazlı ve genel mAP@50", 12, MUTE, italic=True)

# ---------- 9. EDGE / VİZYON ----------
content_slide("Edge / Pazarlanabilirlik Vizyonu", [
    (0, "Hedef: küçük donanımlarda (otobüs/araç) çalışan sistem", TXT, True),
    (0, "Her otobüse: küçük cihaz + kamera + GPS", ACCENT, True),
    (1, "Model yerinde, offline çalışır; çukurları GPS'le işaretler"),
    (1, "Bağlantı gelince merkeze senkronlar (store-and-forward)"),
    (0, "Model küçük (ONNX) → Raspberry Pi / Jetson / telefonda koşar"),
    (0, "Frame sampling + INT8 ile ucuz donanımda gerçek zamanlı", MUTE),
], "Vizyon")

# ---------- 10. TEKNOLOJİ ----------
content_slide("Teknoloji Yığını", [
    (0, "Frontend: Vanilla JS · Canvas/görüntü işleme · ONNX Runtime Web", TXT, True),
    (0, "Yapay Zekâ: YOLOv8 · PyTorch / Ultralytics · ONNX", TXT, True),
    (0, "Backend: Node.js · Express · PostgreSQL", TXT, True),
    (0, "Harita / Panel: Leaflet.js", TXT, True),
    (0, "Eğitim: NVIDIA RTX 5060 (CUDA)", MUTE),
], "Teknoloji")

# ---------- 11. SINIRLAR & GELECEK ----------
content_slide("Sınırlar ve Gelecek Çalışmalar", [
    (0, "Sınır: model yurt dışı verisiyle eğitildi → Türk yolunda recall sınırlı", MUTE, True),
    (0, "Gelecek adımlar:", TXT, True),
    (1, "Türk yolu verisiyle fine-tune → doğruluk artışı (boyut değişmez)"),
    (1, "GPS entegrasyonu + konuma dayalı kümeleme"),
    (1, "Store-and-forward (offline kuyruk)"),
    (1, "INT8 nicemleme ile gömülü cihazda hızlandırma"),
], "Değerlendirme")

# ---------- 12. KAPANIŞ ----------
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, 0, 3.0, 13.333, 1.5, PANEL)
tf = box(s, 1.0, 3.05, 11.3, 1.4); tf.word_wrap = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
set_run(tf.paragraphs[0].add_run(), "Teşekkürler", 40, ACCENT, bold=True)
tf2 = box(s, 1.0, 5.0, 11.3, 1.2)
tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
set_run(tf2.paragraphs[0].add_run(), "Mehmet Emin AKPOLAT  ·  215542003@firat.edu.tr", 18, TXT)
p = tf2.add_paragraph(); p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "github.com/EminAkplt/RoadScan", 15, MUTE)

prs.save(OUT)
print("PPTX olusturuldu:", OUT, "| slayt:", len(prs.slides._sldIdLst))
